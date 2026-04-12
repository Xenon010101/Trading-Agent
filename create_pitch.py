from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_NAVY = RGBColor(10, 25, 41)
ELECTRIC_BLUE = RGBColor(0, 212, 255)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(180, 180, 180)

def add_dark_background(slide, shapes):
    background = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_NAVY
    background.line.fill.background()

def add_title(text_box, text, font_size=44, bold=True):
    text_box.text = text
    pf = text_box.text_frame.paragraphs[0]
    pf.font.size = Pt(font_size)
    pf.font.bold = bold
    pf.font.color.rgb = ELECTRIC_BLUE
    pf.font.name = "Arial Black"

def add_body_text(text_box, text, font_size=18, color=None):
    text_box.text = text
    tf = text_box.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.font.size = Pt(font_size)
        p.font.color.rgb = color if color else WHITE
        p.font.name = "Arial"

def add_slide_number(slide, num, total):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    tf.text = f"{num}/{total}"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

# SLIDE 1 - TITLE
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
add_title(title_box, "InsiderEdge 🤖", 54)

subtitle_box = shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1))
add_body_text(subtitle_box, "Autonomous AI Crypto Trading Agent", 28, ELECTRIC_BLUE)

info_box = shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(2))
add_body_text(info_box, "lablab.ai x Kraken AI Trading Hackathon\n\nTeam: Anmol Patel & Akshita\nGitHub: github.com/Xenon010101/Trading-Agent", 20)

add_slide_number(slide, 1, 12)

# SLIDE 2 - THE PROBLEM
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
add_title(title_box, "The Problem", 40)

content_box = shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(4))
content_box.text_frame.word_wrap = True
p = content_box.text_frame.paragraphs[0]
p.text = "❌ Crypto markets run 24/7 — humans cannot\n❌ Emotional trading leads to losses\n❌ Manual analysis is slow and inconsistent\n❌ No accountability or audit trail"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.font.name = "Arial"

quote_box = shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
add_body_text(quote_box, '"The market is always open. Your emotions are always a liability."', 20, ELECTRIC_BLUE)

add_slide_number(slide, 2, 12)

# SLIDE 3 - THE SOLUTION
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
add_title(title_box, "InsiderEdge", 40)

solution_box = shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(2))
add_body_text(solution_box, "An autonomous AI agent that trades crypto 24/7\nusing real market signals — with zero human emotion", 22, ELECTRIC_BLUE)

checklist = shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(3))
checklist.text_frame.word_wrap = True
p = checklist.text_frame.paragraphs[0]
p.text = "✅ Never sleeps\n✅ Never panics\n✅ Never second-guesses\n✅ Every decision recorded on blockchain"
p.font.size = Pt(24)
p.font.color.rgb = WHITE

add_slide_number(slide, 3, 12)

# SLIDE 4 - HOW IT WORKS
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "How It Works", 36)

flow = shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
flow.text_frame.word_wrap = True
p = flow.text_frame.paragraphs[0]
p.text = "PRISM API → Live market signals\n     ↓\nCoinGecko → Real-time price + OHLC data\n     ↓\nMulti-Signal Scoring Engine\nRSI Score + MACD Score + Trend Score\n     ↓\nGroq LLaMA 70B → Final AI Decision\n     ↓\nRisk Manager → 6-layer safety check\n     ↓\nTrade Executor → Paper / Live trading\n     ↓\nERC-8004 → On-chain verification"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.name = "Consolas"

add_slide_number(slide, 4, 12)

# SLIDE 5 - THE AI BRAIN
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "Groq LLaMA 70B Decision Engine", 36)

steps_box = shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(3))
steps_box.text_frame.word_wrap = True
p = steps_box.text_frame.paragraphs[0]
p.text = "Every 10 minutes InsiderEdge:\n\n1. Fetches live data for BTC, ETH, SOL\n2. Calculates RSI + MACD + Trend scores\n3. Sends data to Groq LLaMA 70B\n4. AI returns: action + confidence + reasoning\n5. Only trades if confidence above 60%"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

example_box = shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2.5))
add_body_text(example_box, "Example Decision:\nAction: BUY\nConfidence: 75%\nReason: MACD bullish + uptrend + positive momentum.\nRSI at 58 — not overbought. Low volatility reduces risk.", 16)

add_slide_number(slide, 5, 12)

# SLIDE 6 - RISK MANAGEMENT
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "6-Layer Risk Management", 36)

rows = 7
cols = 3
table = shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(3)).table

table.cell(0, 0).text = "Layer"
table.cell(0, 1).text = "Rule"
table.cell(0, 2).text = "Value"
for i in range(3):
    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = ELECTRIC_BLUE
    table.cell(0, i).text_frame.paragraphs[0].font.bold = True

data = [
    ("Stop Loss", "Auto-exit losing trade", "0.8%"),
    ("Take Profit", "Auto-lock gains", "1.2%"),
    ("Confidence Gate", "Min AI confidence", "60%"),
    ("Daily Trade Limit", "Max trades per day", "50"),
    ("Circuit Breaker", "Halt if daily loss exceeds", "5%"),
    ("Position Check", "Never sell what you dont own", "Active")
]

for row_idx, (layer, rule, value) in enumerate(data, 1):
    table.cell(row_idx, 0).text = layer
    table.cell(row_idx, 1).text = rule
    table.cell(row_idx, 2).text = value
    for i in range(3):
        table.cell(row_idx, i).text_frame.paragraphs[0].font.color.rgb = WHITE
        table.cell(row_idx, i).text_frame.paragraphs[0].font.size = Pt(14)

circuit_box = shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1))
add_body_text(circuit_box, "Circuit breaker triggered in 4 days: ZERO times", 18, ELECTRIC_BLUE)

add_slide_number(slide, 6, 12)

# SLIDE 7 - ERC-8004 INTEGRATION
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "Trustless On-Chain Verification", 36)

intro_box = shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1.5))
add_body_text(intro_box, "InsiderEdge is registered on Ethereum Sepolia testnet\nvia the ERC-8004 protocol", 18)

info_grid = shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.2))
info_grid.text_frame.word_wrap = True
p = info_grid.text_frame.paragraphs[0]
p.text = "Agent ID: 33     Operator: 0x8cf8...d5F7     Network: Sepolia (Chain ID: 11155111)"
p.font.size = Pt(16)
p.font.color.rgb = ELECTRIC_BLUE
p.font.name = "Consolas"

checks = shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1.5))
checks.text_frame.word_wrap = True
p = checks.text_frame.paragraphs[0]
p.text = "Every trade decision is:\n✅ Signed with EIP-712\n✅ Submitted to RiskRouter contract\n✅ Posted to ValidationRegistry\n✅ Permanently verifiable by anyone"
p.font.size = Pt(16)
p.font.color.rgb = WHITE

add_slide_number(slide, 7, 12)

# SLIDE 8 - RESULTS
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "4 Days of Autonomous Operation", 36)

rows = 9
cols = 2
table = shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(4)).table

table.cell(0, 0).text = "Metric"
table.cell(0, 1).text = "Result"
for i in range(2):
    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = ELECTRIC_BLUE
    table.cell(0, i).text_frame.paragraphs[0].font.bold = True

results_data = [
    ("Total Decisions", "530+"),
    ("Trades Executed", "23+"),
    ("Daily PnL", "+1.60%"),
    ("Circuit Breaker Triggered", "0 times"),
    ("Agent Crashes", "0"),
    ("Uptime", "4 days continuous"),
    ("On-chain Agent ID", "33"),
    ("Blockchain Checkpoints", "Active")
]

for row_idx, (metric, result) in enumerate(results_data, 1):
    table.cell(row_idx, 0).text = metric
    table.cell(row_idx, 1).text = result
    for i in range(2):
        table.cell(row_idx, i).text_frame.paragraphs[0].font.color.rgb = WHITE
        table.cell(row_idx, i).text_frame.paragraphs[0].font.size = Pt(14)

summary_box = shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1))
add_body_text(summary_box, "The agent ran through the night. Every night. Zero supervision.", 18, ELECTRIC_BLUE)

add_slide_number(slide, 8, 12)

# SLIDE 9 - TECH STACK
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "Tech Stack", 36)

rows = 9
cols = 2
table = shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(4)).table

table.cell(0, 0).text = "Layer"
table.cell(0, 1).text = "Technology"
for i in range(2):
    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = ELECTRIC_BLUE
    table.cell(0, i).text_frame.paragraphs[0].font.bold = True

tech_data = [
    ("Language", "Python 3.x"),
    ("AI Brain", "Groq LLaMA 3.3 70B"),
    ("Market Signals", "PRISM API"),
    ("Price Data", "CoinGecko API"),
    ("Blockchain", "Web3.py + Sepolia"),
    ("Exchange", "Kraken CLI"),
    ("Protocol", "ERC-8004"),
    ("Logging", "JSON audit trail")
]

for row_idx, (layer, tech) in enumerate(tech_data, 1):
    table.cell(row_idx, 0).text = layer
    table.cell(row_idx, 1).text = tech
    for i in range(2):
        table.cell(row_idx, i).text_frame.paragraphs[0].font.color.rgb = WHITE
        table.cell(row_idx, i).text_frame.paragraphs[0].font.size = Pt(14)

stats_box = shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
add_body_text(stats_box, "Total lines of code: ~800\nBuilt in: 4 days\nTrading experience before hackathon: Zero", 16, ELECTRIC_BLUE)

add_slide_number(slide, 9, 12)

# SLIDE 10 - PRIZE STRATEGY
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "Prize Strategy", 36)

rows = 5
cols = 3
table = shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(2.5)).table

table.cell(0, 0).text = "Prize"
table.cell(0, 1).text = "Amount"
table.cell(0, 2).text = "Our Edge"
for i in range(3):
    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = ELECTRIC_BLUE
    table.cell(0, i).text_frame.paragraphs[0].font.bold = True

prize_data = [
    ("Best Trustless Agent", "$10,000", "ERC-8004 Agent ID 33"),
    ("Best Risk-Adjusted Return", "$5,000", "6-layer risk system"),
    ("Best Compliance & Risk", "$2,500", "Circuit breaker + audit"),
    ("Social Engagement", "$1,200", "Daily Twitter updates")
]

for row_idx, (prize, amount, edge) in enumerate(prize_data, 1):
    table.cell(row_idx, 0).text = prize
    table.cell(row_idx, 1).text = amount
    table.cell(row_idx, 2).text = edge
    for i in range(3):
        table.cell(row_idx, i).text_frame.paragraphs[0].font.color.rgb = WHITE
        table.cell(row_idx, i).text_frame.paragraphs[0].font.size = Pt(14)

total_box = shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(1))
add_body_text(total_box, "Total targeting: $18,700", 28, ELECTRIC_BLUE)

add_slide_number(slide, 10, 12)

# SLIDE 11 - TEAM
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_title(title_box, "The Team", 36)

anmol_box = shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(2))
anmol_box.text_frame.word_wrap = True
p = anmol_box.text_frame.paragraphs[0]
p.text = "Anmol Patel\n3rd Year Information Technology Student\nRole: Lead Developer — AI brain, trading logic,\nrisk management, ERC-8004 integration\nTwitter: @Anmol_patel2112"
p.font.size = Pt(16)
p.font.color.rgb = WHITE

akshita_box = shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1.5))
akshita_box.text_frame.word_wrap = True
p = akshita_box.text_frame.paragraphs[0]
p.text = "Akshita\n1st Year Computer Science Student\nRole: Co-developer & researcher"
p.font.size = Pt(16)
p.font.color.rgb = WHITE

quote_box = shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(2))
add_body_text(quote_box, "GitHub: github.com/Xenon010101/Trading-Agent\n\nA 1st year and a 3rd year student built a fully\nautonomous on-chain AI trading agent in 4 days.\nZero trading experience. Zero blockchain experience.\nJust determination.", 14, ELECTRIC_BLUE)

add_slide_number(slide, 11, 12)

# SLIDE 12 - THANK YOU
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes
add_dark_background(slide, shapes)

title_box = shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1))
add_title(title_box, "Thank You", 48)

subtitle_box = shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1))
add_title(subtitle_box, "InsiderEdge 🤖", 32)

stats_grid = shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12), Inches(1.2))
stats_grid.text_frame.word_wrap = True
p = stats_grid.text_frame.paragraphs[0]
p.text = "Agent ID: 33 — Sepolia Testnet     Daily PnL: +1.60%     Uptime: 4 days continuous"
p.font.size = Pt(16)
p.font.color.rgb = ELECTRIC_BLUE

contact_box = shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12), Inches(1.2))
contact_box.text_frame.word_wrap = True
p = contact_box.text_frame.paragraphs[0]
p.text = "GitHub: github.com/Xenon010101/Trading-Agent\nTwitter: @Anmol_patel2112\n\nBuilt by Anmol Patel (3rd Year IT) & Akshita (1st Year CSE)\n@krakenfx @lablabai @Surgexyz_"
p.font.size = Pt(14)
p.font.color.rgb = WHITE

add_slide_number(slide, 12, 12)

prs.save("Z:\\Trading Agent\\InsiderEdge_Pitch.pptx")
print("PowerPoint saved as InsiderEdge_Pitch.pptx")